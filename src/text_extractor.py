#!/usr/bin/env python3
"""
Text Extractor (T1.2)

Extracts plain text from non-image artifacts (.md, .txt, .eml, .pdf, .pptx)
and writes one Markdown file per artifact into knowledge/markdown/, with
YAML frontmatter carried over from artifacts/manifest.yaml.

Extracted text is re-scanned for PHI/PII before writing; hits are written
to knowledge/quarantine/ instead and logged.

Usage:
    python3 text_extractor.py <artifact_file> [<artifact_file> ...]
    python3 text_extractor.py --all <artifacts_dir>   # every manifest entry found in dir
"""

import re
import sys
import yaml
from datetime import date
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
from phi_pii_scanner import load_patterns, scan_text  # noqa: E402

ROOT = Path(__file__).parent.parent
MANIFEST = ROOT / "artifacts" / "manifest.yaml"
OUT_DIR = ROOT / "knowledge" / "markdown"
QUARANTINE_DIR = ROOT / "knowledge" / "quarantine"

SUPPORTED = {".md", ".markdown", ".txt", ".eml", ".pdf", ".pptx"}


def clean_text(text):
    """Normalize odd encodings: NEL/CR line endings, control chars, mojibake."""
    text = text.replace("\r\n", "\n").replace("\r", "\n").replace("\x85", "\n")
    # Strip remaining C0/C1 control chars except tab/newline
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]", "", text)
    # Collapse runs of 3+ blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip() + "\n"


def read_with_fallback(path):
    """Try UTF-8 first; fall back to cp1252 (common for Outlook/Windows
    exports), then latin-1 with replacement as last resort."""
    raw = path.read_bytes()
    for enc in ("utf-8", "cp1252"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("latin-1", errors="replace")


def extract_md(path):
    return read_with_fallback(path), "passthrough"


def extract_txt(path):
    return read_with_fallback(path), "plaintext"


def extract_eml(path):
    import email
    from email import policy

    msg = email.message_from_bytes(path.read_bytes(), policy=policy.default)
    parts = [f"Subject: {msg['subject']}", f"From: {msg['from']}", f"Date: {msg['date']}", ""]
    body = msg.get_body(preferencelist=("plain", "html"))
    parts.append(body.get_content() if body else "(no body found)")
    return "\n".join(parts), "email-parser"


def extract_pdf(path):
    from PyPDF2 import PdfReader

    reader = PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        pages.append(f"## Page {i}\n\n{page.extract_text() or '(no extractable text)'}")
    return "\n\n".join(pages), "PyPDF2"


def extract_pptx(path):
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(c.text.strip() for c in row.cells))
        if getattr(slide, "has_notes_slide", False) and slide.notes_slide.notes_text_frame.text.strip():
            texts.append(f"[Speaker notes] {slide.notes_slide.notes_text_frame.text.strip()}")
        body = "\n\n".join(texts) if texts else "(no text on slide — likely diagram/image, needs T1.3 vision extraction)"
        slides.append(f"## Slide {i}\n\n{body}")
    return "\n\n".join(slides), "python-pptx (slide text only; diagrams need T1.3 vision pass)"


EXTRACTORS = {
    ".md": extract_md,
    ".markdown": extract_md,
    ".txt": extract_txt,
    ".eml": extract_eml,
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
}


def load_manifest():
    with open(MANIFEST) as f:
        return yaml.safe_load(f)


def manifest_entry(manifest, filename):
    for a in manifest["artifacts"]:
        if a["filename"] == filename:
            return a
    return None


def slugify(name):
    stem = Path(name).stem.lower()
    return re.sub(r"[^a-z0-9]+", "_", stem).strip("_")


def build_frontmatter(entry, method):
    fm = {
        "source_file": entry["filename"],
        "source": entry["source"],
        "date": entry["date"],
        "state": entry["state"],
        "owner": entry["owner"],
        "category": entry["category"],
        "confidence": entry["confidence"],
        "tags": entry.get("tags", []),
        "extraction_method": method,
        "extracted_on": date.today().isoformat(),
    }
    return "---\n" + yaml.safe_dump(fm, sort_keys=False).strip() + "\n---\n\n"


def extract_file(path, manifest, patterns):
    """Extract one artifact. Returns (status, out_path) where status is
    'ok', 'quarantined', 'skipped', or 'error'."""
    path = Path(path)
    ext = path.suffix.lower()
    if ext not in SUPPORTED:
        print(f"- SKIP (unsupported type {ext}): {path.name}")
        return "skipped", None

    entry = manifest_entry(manifest, path.name)
    if entry is None:
        print(f"- SKIP (not in manifest): {path.name}")
        return "skipped", None

    try:
        raw, method = EXTRACTORS[ext](path)
    except Exception as e:
        print(f"✗ ERROR extracting {path.name}: {e}")
        return "error", None

    text = clean_text(raw)

    # Safety gate: re-scan the extracted text before it enters the KB
    hits = scan_text(text, patterns)
    out_name = slugify(path.name) + ".md"
    content = build_frontmatter(entry, method) + text

    if hits:
        QUARANTINE_DIR.mkdir(parents=True, exist_ok=True)
        out_path = QUARANTINE_DIR / out_name
        out_path.write_text(content)
        print(f"✗ QUARANTINED (PHI in extracted text): {path.name} → {out_path.relative_to(ROOT)}")
        for h in hits:
            print(f"    [{h['severity']}] {h['pattern']} ×{h['count']}: {', '.join(h['masked_samples'])}")
        return "quarantined", out_path

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUT_DIR / out_name
    out_path.write_text(content)
    words = len(text.split())
    print(f"✓ {path.name} → {out_path.relative_to(ROOT)}  ({words} words, via {method})")
    return "ok", out_path


def main():
    args = sys.argv[1:]
    if not args:
        print(__doc__)
        sys.exit(2)

    manifest = load_manifest()
    patterns = load_patterns()

    if args[0] == "--all":
        base = Path(args[1]) if len(args) > 1 else ROOT / "artifacts" / "raw"
        files = [base / a["filename"] for a in manifest["artifacts"]
                 if (base / a["filename"]).exists()]
    else:
        files = [Path(a) for a in args]

    counts = {"ok": 0, "quarantined": 0, "skipped": 0, "error": 0}
    for f in files:
        status, _ = extract_file(f, manifest, patterns)
        counts[status] += 1

    print(f"\nDone: {counts['ok']} extracted, {counts['quarantined']} quarantined, "
          f"{counts['skipped']} skipped, {counts['error']} errors")
    sys.exit(0 if counts["error"] == 0 and counts["quarantined"] == 0 else 1)


if __name__ == "__main__":
    main()
